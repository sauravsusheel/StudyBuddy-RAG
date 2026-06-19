"""
StudyBuddy RAG — Flask backend
Supports: PDF, DOCX, PPTX, PNG, JPG, JPEG, WEBP, GIF, BMP
Pure-Python RAG: file parsers + numpy cosine similarity + OpenAI SDK → OpenRouter
"""

import os
import uuid
import base64
import io
import chromadb
from pathlib import Path

import numpy as np
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()  # loads OPENROUTER_API_KEY from .env if present

app = Flask(__name__, static_folder="static", static_url_path="")
CORS(app)

UPLOAD_FOLDER = Path("uploads")
UPLOAD_FOLDER.mkdir(exist_ok=True)

# In-memory session store for tracking uploaded items in the current browser session
sessions: dict = {}

# Persistent ChromaDB reference database setup
CHROMA_DB_PATH = Path("chroma_db")
CHROMA_DB_PATH.mkdir(exist_ok=True)
chroma_client = chromadb.PersistentClient(path=str(CHROMA_DB_PATH))
collection = chroma_client.get_or_create_collection("compliance_frameworks")

# ── Configuration ─────────────────────────────────────────────────────────────
OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY", "your-openrouter-api-key-here")

OPENROUTER_BASE = "https://openrouter.ai/api/v1"
EMBED_MODEL     = "text-embedding-3-small"
VISION_MODEL    = "openai/gpt-4o"   # used for image OCR/description
CHUNK_SIZE      = 800
CHUNK_OVERLAP   = 150
TOP_K           = 4

ALLOWED_EXTENSIONS = {
    "pdf", "docx", "doc",
    "pptx", "ppt",
    "png", "jpg", "jpeg", "webp", "gif", "bmp"
}

IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp", "gif", "bmp"}


# ── Client ────────────────────────────────────────────────────────────────────

def make_client() -> OpenAI:
    return OpenAI(api_key=OPENROUTER_API_KEY, base_url=OPENROUTER_BASE)


# ── File type detection ───────────────────────────────────────────────────────

def get_ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_image(filename: str) -> bool:
    return get_ext(filename) in IMAGE_EXTENSIONS


# ── Text extractors ───────────────────────────────────────────────────────────

def extract_pdf(path: str) -> list[dict]:
    """Returns [{page, text}, ...]"""
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append({"page": i, "text": text})
    return pages


def extract_docx(path: str) -> list[dict]:
    """Returns [{page, text}, ...] — paragraphs grouped into virtual pages."""
    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    paragraphs.append(t)

    # Group ~20 paragraphs per virtual "page"
    pages = []
    group_size = 20
    for i in range(0, len(paragraphs), group_size):
        text = "\n".join(paragraphs[i : i + group_size])
        pages.append({"page": i // group_size, "text": text})
    return pages


def extract_pptx(path: str) -> list[dict]:
    """Returns [{page, text}, ...] — one entry per slide."""
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        if parts:
            pages.append({"page": i, "text": "\n".join(parts)})
    return pages


def extract_image_via_vision(path: str, filename: str, client: OpenAI) -> list[dict]:
    """
    Send image to GPT-4o vision to extract text / describe content.
    Returns [{page: 0, text: <description>}]
    """
    ext = get_ext(filename)
    mime_map = {
        "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "png": "image/png",  "webp": "image/webp",
        "gif": "image/gif",  "bmp": "image/bmp",
    }
    mime = mime_map.get(ext, "image/png")

    # Resize if very large to keep token cost down (max 1568px on longest side)
    img = Image.open(path)
    max_side = 1568
    if max(img.size) > max_side:
        ratio = max_side / max(img.size)
        new_size = (int(img.width * ratio), int(img.height * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    # Convert to bytes
    buf = io.BytesIO()
    save_fmt = "PNG" if ext in ("png", "bmp", "gif") else "JPEG"
    img.save(buf, format=save_fmt)
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

    response = client.chat.completions.create(
        model=VISION_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        "Please extract ALL text visible in this image exactly as written. "
                        "Then provide a detailed description of any diagrams, charts, or visual content. "
                        "Format: first list extracted text, then describe visuals."
                    )
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime};base64,{b64}"}
                }
            ]
        }],
        max_tokens=2048,
    )
    description = response.choices[0].message.content or ""
    return [{"page": 0, "text": description}]


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_pages(pages: list[dict]) -> list[dict]:
    chunks = []
    for p in pages:
        text = p["text"]
        start = 0
        while start < len(text):
            chunk_text = text[start : start + CHUNK_SIZE].strip()
            if chunk_text:
                chunks.append({
                    "page":     p["page"],
                    "text":     chunk_text,
                    "chunk_id": len(chunks),
                })
            start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


# ── Embeddings & search ───────────────────────────────────────────────────────

def embed_texts(client: OpenAI, texts: list[str]) -> np.ndarray:
    all_vecs = []
    for i in range(0, len(texts), 512):
        batch = texts[i : i + 512]
        resp  = client.embeddings.create(model=EMBED_MODEL, input=batch)
        vecs  = [item.embedding for item in sorted(resp.data, key=lambda x: x.index)]
        all_vecs.extend(vecs)
    arr   = np.array(all_vecs, dtype=np.float32)
    norms = np.linalg.norm(arr, axis=1, keepdims=True)
    norms = np.where(norms == 0, 1, norms)
    return arr / norms


def cosine_top_k(query_vec: np.ndarray, doc_vecs: np.ndarray, k: int) -> list[int]:
    scores = doc_vecs @ query_vec
    return np.argsort(scores)[::-1][: min(k, len(scores))].tolist()


def retrieve_framework_context_for_vendor_text(client: OpenAI, vendor_text: str, session: dict, k: int = 3, max_total_chunks: int = 8) -> list[dict]:
    """
    Extract paragraphs/lines from vendor text, search similarity,
    and return unique matching compliance chunks.
    """
    paragraphs = [p.strip() for p in vendor_text.split("\n\n") if len(p.strip()) > 40]
    if not paragraphs:
        paragraphs = [p.strip() for p in vendor_text.split("\n") if len(p.strip()) > 40]
    if not paragraphs:
        paragraphs = [vendor_text.strip()]

    paragraphs = paragraphs[:15]

    try:
        resp = client.embeddings.create(model=EMBED_MODEL, input=paragraphs)
        para_vecs = [item.embedding for item in sorted(resp.data, key=lambda x: x.index)]
        para_arr = np.array(para_vecs, dtype=np.float32)
        norms = np.linalg.norm(para_arr, axis=1, keepdims=True)
        norms = np.where(norms == 0, 1, norms)
        para_arr = para_arr / norms

        scores = para_arr @ session["embeddings"].T

        top_indices_set = set()
        score_map = {}
        for i in range(len(paragraphs)):
            p_scores = scores[i]
            best_indices = np.argsort(p_scores)[::-1][:k]
            for idx in best_indices:
                top_indices_set.add(idx)
                score_map[idx] = max(score_map.get(idx, 0.0), float(p_scores[idx]))

        sorted_indices = sorted(top_indices_set, key=lambda idx: score_map[idx], reverse=True)
        final_indices = sorted_indices[:max_total_chunks]
        final_indices.sort()

        return [session["chunks"][idx] for idx in final_indices]
    except Exception as e:
        print(f"Error in paragraph retrieval: {e}")
        try:
            q_resp = client.embeddings.create(model=EMBED_MODEL, input=[vendor_text[:4000]])
            q_vec = np.array(q_resp.data[0].embedding, dtype=np.float32)
            q_norm = np.linalg.norm(q_vec)
            if q_norm > 0:
                q_vec /= q_norm
            top_indices = cosine_top_k(q_vec, session["embeddings"], max_total_chunks)
            return [session["chunks"][idx] for idx in top_indices]
        except Exception as f_err:
            print(f"Fallback retrieval failed: {f_err}")
            return session["chunks"][:max_total_chunks]



# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return send_from_directory("static", "index.html")


@app.route("/api/upload", methods=["POST"])
def upload_file():
    if not OPENROUTER_API_KEY or OPENROUTER_API_KEY == "your-openrouter-api-key-here":
        return jsonify({"error": "API key not configured. Add it to .env"}), 500

    if "file" not in request.files:
        return jsonify({"error": "No file in request."}), 400

    file = request.files["file"]
    if not file.filename:
        return jsonify({"error": "No file selected."}), 400

    ext = get_ext(file.filename)
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"Unsupported file type '.{ext}'. Allowed: PDF, DOCX, PPTX, PNG, JPG, JPEG, WEBP, GIF, BMP"}), 400

    session_id = str(uuid.uuid4())
    safe_name  = secure_filename(file.filename)
    file_path  = UPLOAD_FOLDER / f"{session_id}_{safe_name}"
    file.save(str(file_path))

    try:
        client = make_client()

        # ── Extract text based on file type ──
        if ext == "pdf":
            pages = extract_pdf(str(file_path))
            file_type = "PDF"
        elif ext in ("docx", "doc"):
            pages = extract_docx(str(file_path))
            file_type = "Word Document"
        elif ext in ("pptx", "ppt"):
            pages = extract_pptx(str(file_path))
            file_type = "PowerPoint"
        elif ext in IMAGE_EXTENSIONS:
            pages = extract_image_via_vision(str(file_path), safe_name, client)
            file_type = "Image"
        else:
            return jsonify({"error": "Unsupported file type."}), 400

        if not pages:
            return jsonify({"error": "Could not extract any text from this file."}), 400

        # ── Chunk ──
        chunks = chunk_pages(pages)
        if not chunks:
            return jsonify({"error": "File appears to be empty."}), 400

        # ── Embed ──
        texts      = [c["text"] for c in chunks]
        embeddings = embed_texts(client, texts)
        embeddings_list = embeddings.tolist()

        ids = [f"{session_id}_p{c['page'] + 1}_c{idx}" for idx, c in enumerate(chunks)]
        metadatas = [{
            "document_id": safe_name,
            "doc_name_clean": safe_name.replace("_", " ").replace("-", " ").upper(),
            "page_number": c["page"] + 1
        } for c in chunks]

        collection.add(
            ids=ids,
            embeddings=embeddings_list,
            documents=texts,
            metadatas=metadatas
        )

        sessions[session_id] = {
            "file_name":  safe_name,
            "file_type":  file_type,
        }

        label = "slides" if file_type == "PowerPoint" else ("pages" if file_type in ("PDF", "Word Document") else "sections")

        return jsonify({
            "session_id":  session_id,
            "file_name":   safe_name,
            "file_type":   file_type,
            "chunk_count": len(chunks),
            "page_count":  len(pages),
            "message":     f"{file_type} processed and added to reference baseline — {len(pages)} {label}, {len(chunks)} chunks indexed.",
        })

    except Exception as e:
        file_path.unlink(missing_ok=True)
        return jsonify({"error": f"Failed to process file: {str(e)}"}), 500


@app.route("/api/audit", methods=["POST"])
def audit():
    data        = request.get_json(force=True)
    vendor_text = data.get("vendor_text", "").strip()
    model       = data.get("model", "meta-llama/llama-3.1-70b-instruct").strip()

    if not vendor_text:
        return jsonify({"error": "vendor_text is required."}), 400
    if not OPENROUTER_API_KEY or OPENROUTER_API_KEY == "your-openrouter-api-key-here":
        return jsonify({"error": "API key not configured on the server."}), 500

    try:
        client = make_client()

        # 1. Embed query
        q_resp = client.embeddings.create(model=EMBED_MODEL, input=[vendor_text[:4000]])
        query_vector = q_resp.data[0].embedding

        # 2. Query ChromaDB baseline frameworks (NIST, ISO, GDPR, HIPAA, etc.)
        results = collection.query(
            query_embeddings=[query_vector],
            n_results=15
        )

        retrieved_chunks = []
        if results and results["documents"] and len(results["documents"]) > 0:
            docs = results["documents"][0]
            metas = results["metadatas"][0] if results["metadatas"] else []
            ids = results["ids"][0] if results["ids"] else []
            distances = results["distances"][0] if results["distances"] else [0.0] * len(docs)
            
            for i in range(len(docs)):
                meta = metas[i] if i < len(metas) else {}
                dist = distances[i] if i < len(distances) else 0.0
                retrieved_chunks.append({
                    "text": docs[i],
                    "document_id": meta.get("document_id", "Unknown"),
                    "doc_name_clean": meta.get("doc_name_clean", "UNKNOWN FRAMEWORK"),
                    "page_number": meta.get("page_number", 1),
                    "score": dist,
                    "id": ids[i] if i < len(ids) else f"chunk_{i}"
                })

        # --- RETRIEVE AND RERANK PIPELINE PLACEHOLDER ---
        # 1. Deduplicate matches by page unique identifiers: (doc_name_clean, page_number)
        seen_pages = set()
        deduped_chunks = []
        for chunk in retrieved_chunks:
            page_key = (chunk["doc_name_clean"], chunk["page_number"])
            if page_key not in seen_pages:
                seen_pages.add(page_key)
                deduped_chunks.append(chunk)

        # 2. Sort by similarity score (ascending distance in Chroma L2 space)
        # Note: If a reranker like FlashRank is configured:
        # from flashrank import Ranker, RerankRequest
        # ranker = Ranker()
        # rerank_request = RerankRequest(query=vendor_text, passages=[{"id": c["id"], "text": c["text"]} for c in deduped_chunks])
        # reranked_results = ranker.rerank(rerank_request)
        sorted_chunks = sorted(deduped_chunks, key=lambda x: x["score"])

        # 3. Extract the top 5 highest-relevance context pieces
        relevant_chunks = sorted_chunks[:5]

        if not relevant_chunks:
            return jsonify({"error": "No reference compliance frameworks found in the database. Please run ingest.py first."}), 404

        # Build context
        context_parts = []
        for rank, chunk in enumerate(relevant_chunks):
            context_parts.append(
                f"[Doc: {chunk['doc_name_clean']}, Page: {chunk['page_number']}]\n{chunk['text']}"
            )
        context = "\n\n---\n\n".join(context_parts)

        system_prompt = (
            "You are an uncompromising, enterprise-grade Corporate IT Security & Compliance Auditor. "
            "Your sole purpose is to audit a vendor's uploaded policy text against the provided "
            "authoritative corporate compliance frameworks (NIST, ISO, GDPR, and Internal IT Policies)."
        )
        
        user_prompt = (
            "CORE OPERATIONAL CONSTRAINTS:\n"
            "1. TRUTH TO CONTEXT: Rely strictly on the facts directly mentioned in the \"RETRIEVED COMPLIANCE CONTEXT\" section below. If a claim cannot be verified directly by the context, you must explicitly state that it is unverified.\n"
            "2. NO HALLUCINATIONS: Do not assume, extrapolate, or invent internal corporate policies.\n"
            "3. EXPLICIT CITATIONS: Every violation or risk flag you identify MUST be paired with the exact Document Name and Page Number from the retrieved context. Format citations exactly as: [Doc: Name, Page: X].\n\n"
            "REASONING FRAMEWORK (CHAIN OF THOUGHT):\n"
            "Before generating your final audit report, you must systematically process the data through the following mental steps:\n"
            "Step 1: Extract the specific data handling, security, and privacy practices declared in the Vendor's uploaded text.\n"
            "Step 2: Cross-reference those vendor practices against the rules in the retrieved corporate compliance frameworks.\n"
            "Step 3: Check for exact alignment or direct conflicts (e.g., vendor retains data indefinitely vs. internal policy limits retention to 90 days).\n"
            "Step 4: Formulate a final risk verdict based strictly on those conflicts.\n\n"
            "OUTPUT FORMATTING RULES:\n"
            "Your final output must skip all conversational filler and jump directly into a structured markdown report utilizing this exact schema:\n\n"
            "# IT SECURITY & COMPLIANCE AUDIT REPORT\n\n"
            "## 1. EXECUTIVE SUMMARY\n"
            "*   **OVERALL VERDICT:** [APPROVED / CONDITIONAL PASS / FLAGGED & REJECTED]\n"
            "*   **PRIMARY REASON:** [A concise, 1-2 sentence explanation of the highest priority finding]\n\n"
            "## 2. COMPLIANCE VIOLATIONS & RISK FLAGS\n"
            "*(If no violations are found, state \"No compliance conflicts detected.\")*\n"
            "*   **[VIOLATION TYPE / TITLE]**\n"
            "    *   **Vendor Clause:** [Brief description of what the vendor does]\n"
            "    *   **Corporate Policy Conflict:** [Brief description of the corporate rule being broken]\n"
            "    *   **Authoritative Citation:** [Doc: Name, Page: X]\n"
            "    *   **Risk Mitigation Advice:** [Actionable step for the procurement team]\n\n"
            "## 3. DATA TRANSPARENCY & EVIDENCE LOG\n"
            "[A concise markdown table summarizing the audited areas: Audit Area | Vendor Stance | Compliance Status | Citation]\n\n"
            "---\n"
            f"RETRIEVED COMPLIANCE CONTEXT:\n{context}\n\n"
            "---\n"
            f"VENDOR POLICY TEXT TO AUDIT:\n{vendor_text}"
        )

        try:
            completion = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": user_prompt},
                ],
                temperature=0.0,
                max_tokens=1500,
            )
            answer = completion.choices[0].message.content or ""

            # Build citations
            citations = []
            for rank, chunk in enumerate(relevant_chunks):
                snippet = chunk["text"][:300].replace("\n", " ").strip()
                if len(chunk["text"]) > 300:
                    snippet += "..."
                citations.append({
                    "source_num": rank + 1,
                    "page":       chunk["page_number"],
                    "page_label": "Page",
                    "snippet":    snippet,
                    "file_name":  chunk["document_id"],
                    "file_type":  "Framework Document"
                })

            return jsonify({"answer": answer, "citations": citations, "model": model})

        except Exception as api_err:
            print(f"[FAILOVER] OpenRouter LLM Call Failed: {api_err}")
            
            # Construct raw citations fallback
            fallback_citations = []
            for rank, chunk in enumerate(relevant_chunks):
                snippet = chunk["text"][:300].replace("\n", " ").strip()
                if len(chunk["text"]) > 300:
                    snippet += "..."
                fallback_citations.append({
                    "source_num": rank + 1,
                    "page":       chunk["page_number"],
                    "page_label": "Page",
                    "snippet":    snippet,
                    "file_name":  chunk["document_id"],
                    "file_type":  "Framework Document"
                })

            # Return 502 with fallback payload
            return jsonify({
                "is_fallback": True,
                "error": f"Audit Model call failed: {str(api_err)}",
                "message": "The compliance LLM is currently unreachable. Displaying raw matching reference guidelines from the vector database.",
                "citations": fallback_citations
            }), 502

    except Exception as e:
        return jsonify({"error": f"Audit failed: {str(e)}"}), 500


@app.route("/api/chat", methods=["POST"])
def chat():
    data       = request.get_json(force=True)
    question   = data.get("question", "").strip()
    model      = data.get("model", "meta-llama/llama-3.1-8b-instruct").strip()
    mode       = data.get("mode", "chat").strip()

    if not question:
        return jsonify({"error": "question/text is required."}), 400
    if not OPENROUTER_API_KEY or OPENROUTER_API_KEY == "your-openrouter-api-key-here":
        return jsonify({"error": "API key not configured on the server."}), 500

    try:
        client = make_client()

        # 1. Embed query
        q_resp = client.embeddings.create(model=EMBED_MODEL, input=[question])
        query_vector = q_resp.data[0].embedding

        # 2. Query ChromaDB
        results = collection.query(
            query_embeddings=[query_vector],
            n_results=15 if mode == "audit" else TOP_K
        )

        retrieved_chunks = []
        if results and results["documents"] and len(results["documents"]) > 0:
            docs = results["documents"][0]
            metas = results["metadatas"][0] if results["metadatas"] else []
            ids = results["ids"][0] if results["ids"] else []
            distances = results["distances"][0] if results["distances"] else [0.0] * len(docs)
            
            for i in range(len(docs)):
                meta = metas[i] if i < len(metas) else {}
                dist = distances[i] if i < len(distances) else 0.0
                retrieved_chunks.append({
                    "text": docs[i],
                    "document_id": meta.get("document_id", "Unknown"),
                    "doc_name_clean": meta.get("doc_name_clean", "UNKNOWN FRAMEWORK"),
                    "page_number": meta.get("page_number", 1),
                    "score": dist,
                    "id": ids[i] if i < len(ids) else f"chunk_{i}"
                })

        # --- RETRIEVE AND RERANK PIPELINE PLACEHOLDER ---
        seen_pages = set()
        deduped_chunks = []
        for chunk in retrieved_chunks:
            page_key = (chunk["doc_name_clean"], chunk["page_number"])
            if page_key not in seen_pages:
                seen_pages.add(page_key)
                deduped_chunks.append(chunk)

        sorted_chunks = sorted(deduped_chunks, key=lambda x: x["score"])
        relevant_chunks = sorted_chunks[:5 if mode == "audit" else TOP_K]

        if not relevant_chunks:
            return jsonify({"error": "No reference compliance frameworks found in the database. Please run ingest.py first."}), 404

        context_parts = []
        for rank, chunk in enumerate(relevant_chunks):
            context_parts.append(
                f"[Doc: {chunk['doc_name_clean']}, Page: {chunk['page_number']}]\n{chunk['text']}"
            )
        context = "\n\n---\n\n".join(context_parts)

        if mode == "audit":
            system_prompt = (
                "You are an uncompromising, enterprise-grade Corporate IT Security & Compliance Auditor. "
                "Your sole purpose is to audit a vendor's uploaded policy text against the provided "
                "authoritative corporate compliance frameworks (NIST, ISO, GDPR, and Internal IT Policies)."
            )
            
            user_prompt = (
                "CORE OPERATIONAL CONSTRAINTS:\n"
                "1. TRUTH TO CONTEXT: Rely strictly on the facts directly mentioned in the \"RETRIEVED COMPLIANCE CONTEXT\" section below. If a claim cannot be verified directly by the context, you must explicitly state that it is unverified.\n"
                "2. NO HALLUCINATIONS: Do not assume, extrapolate, or invent internal corporate policies.\n"
                "3. EXPLICIT CITATIONS: Every violation or risk flag you identify MUST be paired with the exact Document Name and Page Number from the retrieved context. Format citations exactly as: [Doc: Name, Page: X].\n\n"
                "REASONING FRAMEWORK (CHAIN OF THOUGHT):\n"
                "Before generating your final audit report, you must systematically process the data through the following mental steps:\n"
                "Step 1: Extract the specific data handling, security, and privacy practices declared in the Vendor's uploaded text.\n"
                "Step 2: Cross-reference those vendor practices against the rules in the retrieved corporate compliance frameworks.\n"
                "Step 3: Check for exact alignment or direct conflicts (e.g., vendor retains data indefinitely vs. internal policy limits retention to 90 days).\n"
                "Step 4: Formulate a final risk verdict based strictly on those conflicts.\n\n"
                "OUTPUT FORMATTING RULES:\n"
                "Your final output must skip all conversational filler and jump directly into a structured markdown report utilizing this exact schema:\n\n"
                "# IT SECURITY & COMPLIANCE AUDIT REPORT\n\n"
                "## 1. EXECUTIVE SUMMARY\n"
                "*   **OVERALL VERDICT:** [APPROVED / CONDITIONAL PASS / FLAGGED & REJECTED]\n"
                "*   **PRIMARY REASON:** [A concise, 1-2 sentence explanation of the highest priority finding]\n\n"
                "## 2. COMPLIANCE VIOLATIONS & RISK FLAGS\n"
                "*(If no violations are found, state \"No compliance conflicts detected.\")*\n"
                "*   **[VIOLATION TYPE / TITLE]**\n"
                "    *   **Vendor Clause:** [Brief description of what the vendor does]\n"
                "    *   **Corporate Policy Conflict:** [Brief description of the corporate rule being broken]\n"
                "    *   **Authoritative Citation:** [Doc: Name, Page: X]\n"
                "    *   **Risk Mitigation Advice:** [Actionable step for the procurement team]\n\n"
                "## 3. DATA TRANSPARENCY & EVIDENCE LOG\n"
                "[A concise markdown table summarizing the audited areas: Audit Area | Vendor Stance | Compliance Status | Citation]\n\n"
                "---\n"
                f"RETRIEVED COMPLIANCE CONTEXT:\n{context}\n\n"
                "---\n"
                f"VENDOR POLICY TEXT TO AUDIT:\n{question}"
            )
            temperature = 0.0
            max_tokens = 1500
        else:
            system_prompt = (
                "You are SecurAudit, a corporate IT security & compliance assistant. "
                "Answer the user's questions about the uploaded security and compliance frameworks "
                "using ONLY the provided context. If the answer is not in the context, state that it is unverified."
            )
            user_prompt = (
                f"Context from compliance frameworks:\n\n{context}\n\n"
                f"Question: {question}\n\nAnswer:"
            )
            temperature = 0.2
            max_tokens = 1024

        try:
            completion = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": user_prompt},
                ],
                temperature=temperature,
                max_tokens=max_tokens,
            )
            answer = completion.choices[0].message.content or ""

            citations = []
            for rank, chunk in enumerate(relevant_chunks):
                snippet = chunk["text"][:300].replace("\n", " ").strip()
                if len(chunk["text"]) > 300:
                    snippet += "..."
                citations.append({
                    "source_num": rank + 1,
                    "page":       chunk["page_number"],
                    "page_label": "Page",
                    "snippet":    snippet,
                    "file_name":  chunk["document_id"],
                    "file_type":  "Framework Document"
                })

            return jsonify({"answer": answer, "citations": citations, "model": model})

        except Exception as api_err:
            if mode == "audit":
                print(f"[FAILOVER] OpenRouter API Failed: {api_err}")
                fallback_citations = []
                for rank, chunk in enumerate(relevant_chunks):
                    snippet = chunk["text"][:300].replace("\n", " ").strip()
                    if len(chunk["text"]) > 300:
                        snippet += "..."
                    fallback_citations.append({
                        "source_num": rank + 1,
                        "page":       chunk["page_number"],
                        "page_label": "Page",
                        "snippet":    snippet,
                        "file_name":  chunk["document_id"],
                        "file_type":  "Framework Document"
                    })
                return jsonify({
                    "is_fallback": True,
                    "error": f"LLM API Execution Failed: {str(api_err)}",
                    "message": "Unable to contact the audit model. Displaying raw compliance reference clauses retrieved from the database.",
                    "citations": fallback_citations
                }), 502
            else:
                raise api_err

    except Exception as e:
        return jsonify({"error": f"Request failed: {str(e)}"}), 500


@app.route("/api/session/<session_id>", methods=["DELETE"])
def delete_session(session_id):
    sessions.pop(session_id, None)
    for f in UPLOAD_FOLDER.glob(f"{session_id}_*"):
        f.unlink(missing_ok=True)
    return jsonify({"message": "Session deleted."})


if __name__ == "__main__":
    app.run(debug=True, port=5000)
